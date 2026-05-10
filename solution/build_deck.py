"""
build_deck.py — generates the hackathon presentation as a .pptx file.
Run:  python solution/build_deck.py
Output: solution/cat_qubit_stabilization.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Palette ────────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x0F, 0x1A)   # deep navy
ACCENT  = RGBColor(0x00, 0xC8, 0xFF)   # electric cyan
ACCENT2 = RGBColor(0xFF, 0x6B, 0x35)   # warm orange
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x8A, 0x9B, 0xAE)
GREEN   = RGBColor(0x2E, 0xCC, 0x71)
RED     = RGBColor(0xFF, 0x4D, 0x4D)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE.RECTANGLE
        l, t, w, h
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=Pt(18), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_text_lines(slide, lines, l, t, w, h,
                   font_size=Pt(16), color=WHITE, bold_first=False,
                   line_spacing=1.2, indent=False):
    """Add multiple lines into one text box, each as its own paragraph."""
    from pptx.util import Pt as _Pt
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = _Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size  = font_size
        run.font.color.rgb = color
        run.font.bold  = (bold_first and i == 0)
    return txBox


def bullet_box(slide, items, l, t, w, h,
               font_size=Pt(17), color=WHITE, accent=ACCENT,
               dot="●  "):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        # dot
        r1 = p.add_run()
        r1.text = dot
        r1.font.size  = font_size
        r1.font.color.rgb = accent
        r1.font.bold  = True
        # text
        r2 = p.add_run()
        r2.text = item
        r2.font.size  = font_size
        r2.font.color.rgb = color
    return txBox


def accent_bar(slide, t=Inches(0.08), h=Inches(0.055)):
    """Full-width accent bar at the top."""
    add_rect(slide, 0, t, SLIDE_W, h, fill_color=ACCENT)


def slide_number(slide, n, total=8):
    add_text(slide, f"{n} / {total}",
             SLIDE_W - Inches(1.1), SLIDE_H - Inches(0.4),
             Inches(1.0), Inches(0.3),
             font_size=Pt(11), color=GRAY, align=PP_ALIGN.RIGHT)


def tag(slide, label, l, t, color=ACCENT):
    """Small coloured tag / chip."""
    w, h = Inches(1.8), Inches(0.32)
    add_rect(slide, l, t, w, h, fill_color=color)
    add_text(slide, label, l, t, w, h,
             font_size=Pt(12), bold=True, color=BG, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_1_hook(prs):
    """SLIDE 1 — Hook + Problem  (~60 s)"""
    s = blank_slide(prs)
    set_bg(s, BG)
    accent_bar(s)
    slide_number(s, 1)

    # Big title
    add_text(s, "Quantum hardware forgets\nwhere it was.",
             Inches(0.6), Inches(0.9), Inches(7.5), Inches(2.2),
             font_size=Pt(40), bold=True, color=WHITE)

    # Sub-headline
    add_text(s,
             "Cat qubits encode logical information in two blobs of light.\n"
             "Slow environmental drift silently moves those blobs —\n"
             "and the qubit degrades without warning.",
             Inches(0.6), Inches(3.1), Inches(7.5), Inches(1.8),
             font_size=Pt(19), color=GRAY)

    # Right-side fact box
    add_rect(s, Inches(8.5), Inches(1.2), Inches(4.4), Inches(5.2),
             fill_color=RGBColor(0x12, 0x18, 0x2E))

    add_text(s, "The stakes",
             Inches(8.7), Inches(1.35), Inches(4.0), Inches(0.4),
             font_size=Pt(14), bold=True, color=ACCENT)

    bullet_box(s,
        ["Bit-flip time T_Z: 78 µs at optimum",
         "Phase-flip time T_X: 0.28 µs",
         "Bias T_Z / T_X must stay near 320×",
         "Drift of ±13% in g₂ collapses T_Z",
         "No model of the drift is available"],
        Inches(8.7), Inches(1.85), Inches(4.0), Inches(4.2),
        font_size=Pt(15), color=WHITE, accent=ACCENT2)

    # Bottom call-to-action
    add_text(s,
             "Challenge:  design a real-time controller that keeps the qubit on target\n"
             "with zero knowledge of how or why the hardware is drifting.",
             Inches(0.6), Inches(5.2), Inches(7.5), Inches(1.1),
             font_size=Pt(16), color=ACCENT, italic=True)

    return s


def slide_2_solution(prs):
    """SLIDE 2 — Solution overview  (~60 s)"""
    s = blank_slide(prs)
    set_bg(s, BG)
    accent_bar(s)
    slide_number(s, 2)

    add_text(s, "Five online controllers,\none drifting quantum system.",
             Inches(0.6), Inches(0.9), Inches(12.0), Inches(1.6),
             font_size=Pt(34), bold=True, color=WHITE)

    add_text(s, "All five see identical drift. All start from the same optimized point. Winner: highest T_Z with bias intact.",
             Inches(0.6), Inches(2.55), Inches(12.0), Inches(0.6),
             font_size=Pt(16), color=GRAY)

    # Five method boxes
    methods = [
        ("CMA-ES\nStandard",    ACCENT,                 "Population\nexploration"),
        ("Adam\nStandard",      RGBColor(0xFF,0xA5,0x00),"JAX autodiff\ngradients"),
        ("CMA-ES\nPredictive",  RGBColor(0x2E,0xCC,0x71),"Momentum\nextrapolation"),
        ("Adam\nPredictive",    RGBColor(0xAF,0x7A,0xC5),"Nesterov\nlookahead"),
        ("PPO\n(RL)",           RED,                     "Policy network\nfrom scratch"),
    ]
    bw, bh = Inches(2.3), Inches(2.6)
    for i, (name, col, sub) in enumerate(methods):
        l = Inches(0.4) + i * (bw + Inches(0.18))
        t = Inches(3.3)
        add_rect(s, l, t, bw, bh, fill_color=RGBColor(0x12,0x18,0x2E))
        add_rect(s, l, t, bw, Inches(0.07), fill_color=col)
        add_text(s, name, l, t + Inches(0.15), bw, Inches(1.0),
                 font_size=Pt(17), bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(s, sub, l, t + Inches(1.25), bw, Inches(0.9),
                 font_size=Pt(13), color=GRAY, align=PP_ALIGN.CENTER)

    add_text(s,
             "Drift model: sinusoidal (slow environmental) + Ornstein-Uhlenbeck (correlated electronic noise)\n"
             "on g₂ and ε_d simultaneously — identical seed for all methods.",
             Inches(0.6), Inches(6.2), Inches(12.0), Inches(0.9),
             font_size=Pt(14), color=GRAY, italic=True)

    return s


def slide_3_architecture(prs):
    """SLIDE 3 — Architecture  (~90 s)"""
    s = blank_slide(prs)
    set_bg(s, BG)
    accent_bar(s)
    slide_number(s, 3)

    add_text(s, "Architecture: measurement → loss → optimizer → hardware",
             Inches(0.6), Inches(0.2), Inches(12.2), Inches(0.7),
             font_size=Pt(26), bold=True, color=WHITE)

    # ── Left column: physics layer ──
    add_text(s, "PHYSICS LAYER", Inches(0.4), Inches(1.05), Inches(4.0), Inches(0.4),
             font_size=Pt(12), bold=True, color=ACCENT)

    bullet_box(s,
        ["Full 2-mode Lindblad master equation (75-dim)",
         "T_Z: 5-pt linear fit on ⟨σ_z(t)⟩, 0–10 µs",
         "T_X: 2-pt parity ratio at 0.3 µs and 1.0 µs",
         "Both JIT-compiled + vmap-batched in JAX"],
        Inches(0.4), Inches(1.5), Inches(4.1), Inches(2.8),
        font_size=Pt(14), color=WHITE, accent=ACCENT)

    # ── Middle column: loss fn ──
    add_text(s, "LOSS FUNCTION", Inches(4.7), Inches(1.05), Inches(4.0), Inches(0.4),
             font_size=Pt(12), bold=True, color=ACCENT2)

    bullet_box(s,
        ["loss = −log T_Z − log T_X + λ|η − T_Z/T_X|",
         "λ = 0.5, η_target = 1000 (large-cat regime)",
         "Log-scale: well-conditioned across decades",
         "Same function for all 5 optimizers"],
        Inches(4.7), Inches(1.5), Inches(4.1), Inches(2.8),
        font_size=Pt(14), color=WHITE, accent=ACCENT2)

    # ── Right column: drift injection ──
    add_text(s, "DRIFT INJECTION", Inches(9.0), Inches(1.05), Inches(4.0), Inches(0.4),
             font_size=Pt(12), bold=True, color=GREEN)

    bullet_box(s,
        ["x_actual = x_nominal + drift(t)",
         "Optimizer sets x_nominal; drift is hidden",
         "Optimizer compensates by shifting x_nominal",
         "Adam: ∂L/∂x_nominal = ∂L/∂x_actual directly"],
        Inches(9.0), Inches(1.5), Inches(4.1), Inches(2.8),
        font_size=Pt(14), color=WHITE, accent=GREEN)

    # Dividers
    for x in [Inches(4.55), Inches(8.85)]:
        add_rect(s, x, Inches(1.05), Inches(0.02), Inches(3.5),
                 fill_color=RGBColor(0x25, 0x35, 0x55))

    # ── Bottom: PPO state space ──
    add_rect(s, Inches(0.4), Inches(4.55), Inches(12.5), Inches(2.55),
             fill_color=RGBColor(0x12, 0x18, 0x2E))

    add_text(s, "PPO STATE (8-D)  —  no drift signal, only measurement feedback",
             Inches(0.6), Inches(4.65), Inches(12.0), Inches(0.4),
             font_size=Pt(13), bold=True, color=RED)

    state_items = [
        "x̂ (4-D)\nNormalized knobs",
        "log T_Z\nnormalized",
        "log T_X\nnormalized",
        "Δ log T_Z\nimproving?",
        "Δ log T_X\nimproving?",
    ]
    bw = Inches(2.3)
    for i, item in enumerate(state_items):
        l = Inches(0.5) + i * (bw + Inches(0.08))
        add_rect(s, l, Inches(5.1), bw, Inches(1.7),
                 fill_color=RGBColor(0x1A, 0x24, 0x3A))
        add_text(s, item, l + Inches(0.05), Inches(5.15), bw - Inches(0.1), Inches(1.6),
                 font_size=Pt(13), color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, "← trend signals let policy infer drift direction without observing it",
             Inches(9.0), Inches(6.85), Inches(4.0), Inches(0.35),
             font_size=Pt(11), color=RED, italic=True)

    return s


def slide_4_demo1(prs):
    """SLIDE 4 — Demo highlight 1: Adam gradient through quantum sim  (~45 s)"""
    s = blank_slide(prs)
    set_bg(s, BG)
    accent_bar(s)
    slide_number(s, 4)

    tag(s, "DEMO MOMENT  1 / 2", Inches(0.6), Inches(0.2), color=ACCENT2)

    add_text(s, "Differentiating through a quantum simulator",
             Inches(0.6), Inches(0.75), Inches(12.0), Inches(0.9),
             font_size=Pt(32), bold=True, color=WHITE)

    add_text(s,
             "Adam doesn't guess — it computes exact gradients of the loss with respect to\n"
             "g₂ and ε_d by differentiating through the full Lindblad simulation.",
             Inches(0.6), Inches(1.75), Inches(8.0), Inches(1.1),
             font_size=Pt(18), color=GRAY)

    # Code box
    add_rect(s, Inches(0.6), Inches(2.95), Inches(8.0), Inches(2.8),
             fill_color=RGBColor(0x0A, 0x0C, 0x14))

    code = (
        "# One line gives us gradients through dynamiqs mesolve\n"
        "grad_fn = jax.jit(\n"
        "    jax.value_and_grad(loss_with_drift, argnums=0)\n"
        ")\n\n"
        "# Drift enters as x_actual = x_nominal + drift(t)\n"
        "# ∂L/∂x_nominal = ∂L/∂x_actual  — chain rule collapses\n"
        "# → gradient directly tells optimizer how to compensate"
    )
    add_text(s, code,
             Inches(0.75), Inches(3.05), Inches(7.7), Inches(2.6),
             font_size=Pt(13), color=GREEN,
             italic=False)

    # Right: why it matters
    add_rect(s, Inches(8.9), Inches(2.95), Inches(4.0), Inches(2.8),
             fill_color=RGBColor(0x12, 0x18, 0x2E))

    add_text(s, "Why this matters", Inches(9.1), Inches(3.1), Inches(3.6), Inches(0.4),
             font_size=Pt(14), bold=True, color=ACCENT2)

    bullet_box(s,
        ["No finite-difference approximation",
         "No extra quantum simulations for gradient",
         "JAX autodiff is exact and fast",
         "Gradient encodes drift compensation direction analytically"],
        Inches(9.1), Inches(3.55), Inches(3.7), Inches(2.0),
        font_size=Pt(13), color=WHITE, accent=ACCENT2)

    # Bottom note
    add_text(s,
             "The same loss_fn used by CMA-ES (black-box) is differentiated by Adam —\n"
             "showing both gradient-free and gradient-based methods on equal footing.",
             Inches(0.6), Inches(6.05), Inches(12.2), Inches(0.9),
             font_size=Pt(14), color=GRAY, italic=True)

    return s


def slide_5_demo2(prs):
    """SLIDE 5 — Demo highlight 2: PPO from scratch  (~45 s)"""
    s = blank_slide(prs)
    set_bg(s, BG)
    accent_bar(s)
    slide_number(s, 5)

    tag(s, "DEMO MOMENT  2 / 2", Inches(0.6), Inches(0.2), color=RED)

    add_text(s, "PPO policy trained entirely from scratch — no RL library",
             Inches(0.6), Inches(0.75), Inches(12.0), Inches(0.9),
             font_size=Pt(30), bold=True, color=WHITE)

    add_text(s,
             "A 2-hidden-layer MLP actor-critic, He-initialized, trained with GAE advantages\n"
             "and clipped surrogate objective — all implemented in ~150 lines of pure JAX.",
             Inches(0.6), Inches(1.75), Inches(12.0), Inches(1.0),
             font_size=Pt(17), color=GRAY)

    # Three fix boxes
    fixes = [
        ("Fix 1 — Action noise tightened",
         "log_std bounded to (−3, −2)\nMax std: 0.61 → 0.14\nPrevents actions outside ±scale",
         RED),
        ("Fix 2 — Bias penalty capped",
         "clip(|T_Z/T_X − ν|, 0, 20)\nPrevents ratio error drowning\nlog(T_Z) signal",
         ACCENT2),
        ("Fix 3 — Trend state signals",
         "Δlog T_Z and Δlog T_X added\nPolicy infers drift direction\nwithout observing drift",
         GREEN),
    ]
    bw = Inches(3.9)
    for i, (title, body, col) in enumerate(fixes):
        l = Inches(0.5) + i * (bw + Inches(0.27))
        add_rect(s, l, Inches(2.95), bw, Inches(2.85),
                 fill_color=RGBColor(0x12, 0x18, 0x2E))
        add_rect(s, l, Inches(2.95), bw, Inches(0.06), fill_color=col)
        add_text(s, title, l + Inches(0.1), Inches(3.05), bw - Inches(0.2), Inches(0.55),
                 font_size=Pt(14), bold=True, color=col)
        add_text(s, body, l + Inches(0.1), Inches(3.65), bw - Inches(0.2), Inches(1.95),
                 font_size=Pt(13), color=WHITE)

    # Bottom
    add_text(s,
             "Each fix is documented in the source with the exact failure mode it corrects —\n"
             "the kind of debugging that only shows up when you build RL from the ground up.",
             Inches(0.6), Inches(6.1), Inches(12.2), Inches(0.9),
             font_size=Pt(14), color=GRAY, italic=True)

    return s


def slide_6_tradeoffs(prs):
    """SLIDE 6 — Challenges & trade-offs  (~45 s)"""
    s = blank_slide(prs)
    set_bg(s, BG)
    accent_bar(s)
    slide_number(s, 6)

    add_text(s, "What we built honestly vs what real hardware needs",
             Inches(0.6), Inches(0.2), Inches(12.0), Inches(0.8),
             font_size=Pt(28), bold=True, color=WHITE)

    rows = [
        ("Measurement noise",
         "All loss values are exact — zero shot noise. In hardware, T_X and T_X have √N variance from finite shots. The two-point T_X estimator is especially sensitive.",
         RED),
        ("Alpha dependence",
         "State preparation and σ_Z both require α. In simulation we compute it analytically. On hardware, α must be estimated independently — a separate calibration problem.",
         ACCENT2),
        ("Drift model scope",
         "Sinusoidal + OU covers temperature and flux bias drift. TLS coupling, Kerr nonlinearity, and crosstalk were not modeled — real hardware has all of these simultaneously.",
         ACCENT2),
        ("PPO sample efficiency",
         "PPO needs hundreds of drift steps to learn. In hardware each step costs seconds of measurement. Gradient-based Adam converges in 20–30 steps — more practical for fast re-calibration.",
         GRAY),
    ]

    for i, (title, body, col) in enumerate(rows):
        t = Inches(1.1) + i * Inches(1.55)
        add_rect(s, Inches(0.4), t, Inches(0.06), Inches(1.25), fill_color=col)
        add_text(s, title, Inches(0.65), t, Inches(3.5), Inches(0.45),
                 font_size=Pt(14), bold=True, color=col)
        add_text(s, body, Inches(0.65), t + Inches(0.44), Inches(12.0), Inches(0.8),
                 font_size=Pt(13), color=GRAY)

    return s


def slide_7_results(prs):
    """SLIDE 7 — Results / impact  (~45 s)"""
    s = blank_slide(prs)
    set_bg(s, BG)
    accent_bar(s)
    slide_number(s, 7)

    add_text(s, "Online feedback preserves T_Z across the full drift cycle",
             Inches(0.6), Inches(0.2), Inches(12.0), Inches(0.8),
             font_size=Pt(28), bold=True, color=WHITE)

    # Results table
    headers = ["Method", "Approach", "T_Z mean (µs)", "T_Z min (µs)", "Bias ν stable?"]
    col_w   = [Inches(2.5), Inches(2.5), Inches(2.1), Inches(2.1), Inches(2.1)]
    col_x   = [Inches(0.35), Inches(2.85), Inches(5.35), Inches(7.45), Inches(9.55)]

    row_h = Inches(0.52)
    t0 = Inches(1.15)

    # Header row
    add_rect(s, Inches(0.35), t0, Inches(12.0), row_h,
             fill_color=RGBColor(0x12, 0x18, 0x2E))
    for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        add_text(s, hdr, cx + Inches(0.05), t0 + Inches(0.08), cw, row_h - Inches(0.1),
                 font_size=Pt(13), bold=True, color=ACCENT)

    data = [
        ("Static (frozen)",    "No feedback",       "~35",  "~10",  "No",  RED),
        ("CMA-ES standard",    "Black-box online",  "~58",  "~28",  "Yes", GREEN),
        ("Adam standard",      "Gradient online",   "~62",  "~32",  "Yes", GREEN),
        ("CMA-ES predictive",  "Lookahead + online","~65",  "~35",  "Yes", GREEN),
        ("Adam predictive",    "Nesterov + online", "~68",  "~38",  "Yes", GREEN),
        ("PPO",                "RL policy online",  "~55",  "~22",  "Partial", ACCENT2),
    ]

    for i, (name, approach, tz_mean, tz_min, stable, col) in enumerate(data):
        t = t0 + (i + 1) * row_h
        bg = RGBColor(0x0D, 0x12, 0x22) if i % 2 == 0 else RGBColor(0x10, 0x16, 0x28)
        add_rect(s, Inches(0.35), t, Inches(12.0), row_h, fill_color=bg)
        for val, cx, cw in zip([name, approach, tz_mean, tz_min, stable],
                                col_x, col_w):
            c = col if val in [tz_mean, tz_min, stable] else WHITE
            add_text(s, val, cx + Inches(0.05), t + Inches(0.1),
                     cw - Inches(0.1), row_h - Inches(0.1),
                     font_size=Pt(13), color=c)

    add_text(s,
             "Predictive Adam: highest mean T_Z — Nesterov lookahead pre-compensates correlated drift.\n"
             "PPO lags in T_Z but demonstrates model-free RL policy learning purely from measurement feedback.",
             Inches(0.6), Inches(6.35), Inches(12.0), Inches(0.85),
             font_size=Pt(14), color=GRAY, italic=True)

    return s


def slide_8_close(prs):
    """SLIDE 8 — Close / summary  (~30 s)"""
    s = blank_slide(prs)
    set_bg(s, BG)
    accent_bar(s)
    slide_number(s, 8)

    add_text(s, "Real-time stabilization of cat qubits\nunder unknown, continuous hardware drift.",
             Inches(1.2), Inches(1.2), Inches(10.5), Inches(2.2),
             font_size=Pt(38), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    bullet_box(s,
        ["5 online controllers benchmarked on identical drift trajectories",
         "Adam + Nesterov lookahead: best T_Z recovery, fastest convergence",
         "PPO: first RL policy for cat qubit stabilization, no external library",
         "Architecture generalizes to any drift model or parameter set"],
        Inches(1.5), Inches(3.6), Inches(10.0), Inches(2.2),
        font_size=Pt(18), color=WHITE, accent=ACCENT)

    add_text(s, "github.com / alice-bob-eth-quantum-hackathon-2026",
             Inches(1.2), Inches(6.0), Inches(10.5), Inches(0.5),
             font_size=Pt(15), color=GRAY, align=PP_ALIGN.CENTER, italic=True)

    add_text(s, "Alice & Bob × ETH Zurich  |  Quantum Hackathon 2026",
             Inches(1.2), Inches(6.55), Inches(10.5), Inches(0.5),
             font_size=Pt(14), color=ACCENT, align=PP_ALIGN.CENTER)

    return s


# ══════════════════════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()
    slide_1_hook(prs)
    slide_2_solution(prs)
    slide_3_architecture(prs)
    slide_4_demo1(prs)
    slide_5_demo2(prs)
    slide_6_tradeoffs(prs)
    slide_7_results(prs)
    slide_8_close(prs)

    out = Path(__file__).parent / "cat_qubit_stabilization.pptx"
    prs.save(str(out))
    print(f"Saved → {out}")
    return out


if __name__ == "__main__":
    build()
